import streamlit as st
import ast
from pptx.util import Inches, Pt
from pptx import Presentation
import base64
import openai
import re
from collections.abc import Iterable

MAX_TOKENS = 4096  # Maximum tokens allowed in a single API call
TOKENS_PER_SLIDE_ESTIMATE = 200  # Rough estimate of tokens used per slide

# Set OpenAI API key
openai.api_key = st.secrets['OPENAI_KEY']

def generate_slide_content(title, engine='gpt-3.5-turbo'):
    slide_content = {}
    prompts = [
        ("One descriptive short title of 5-7 words for this slide", 20),
        ("Three useful bullets of 10-14 words each", 50),
        ("One short key takeaway message of 8 words or less", 20),
        ("Five detailed talking points of 30-40 words each", 100)
    ]
    
    key_mapping = {
        "One descriptive short title of 5-7 words for this slide": "crisp_title",
        "Three useful bullets of 10-14 words each": "bullets",
        "One short key takeaway message of 8 words or less": "takeaway_message",
        "Five detailed talking points of 30-40 words each": "talking_points"
    }

    for prompt, max_tokens in prompts:
        response = openai.ChatCompletion.create(
            model=engine,
            messages=[
                {"role": "system", "content": "You are a helpful assistant."},
                {"role": "user", "content": f"Generate content for the title: '{title}'\n\n Create {prompt}:\n1."}
            ]
        )

        result = response['choices'][0]['message']['content'].strip().split('\n')
        result = [r.strip() for r in result if r.strip()]
        
        # Remove any leading numbers from the AI's output to prevent double numbering
        result = [re.sub(r'^\d+\.\s*', '', r) for r in result]

        #slide_content[prompt.lower().replace(" ", "_")] = result[0] if len(result) == 1 else result
        slide_content[key_mapping[prompt]] = result[0] if len(result) == 1 else result


    return slide_content

def generate_outline(presentation_topic, num_slides, engine='gpt-3.5-turbo'):
    response = openai.ChatCompletion.create(
        model=engine,
        messages=[
            {"role": "system", "content": "You are a helpful assistant."},
            {"role": "user", "content": f"Generate {num_slides} slide titles for a presentation on the topic: '{presentation_topic}'."}
        ]
    )

    generated_text = response['choices'][0]['message']['content'].strip()

    outline = generated_text.split('\n')
    outline = [slide[slide.find(":")+1:].strip() if ":" in slide else slide.strip() for slide in outline]

    if len(outline) != num_slides:
        print(f"Warning: Expected {num_slides} slide titles but received {len(outline)}")

    return outline

def create_presentation(slides_content, company_name, presentation_name, presenter):
    # Initialize a Presentation object
    presentation = Presentation()

    # Add title slide
    slide_layout = presentation.slide_layouts[0]  # Slide layout 0 is a title slide
    slide = presentation.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = presentation_name
    subtitle = slide.placeholders[1]
    subtitle.text = presenter

    # Add other slides
    for slide_content in slides_content:
        # Add a new slide with a title and content layout
        slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(slide_layout)

        # Set the title
        title = slide.shapes.title
        title.text = slide_content["crisp_title"]

        # Add bullet points
        content = slide.placeholders[1]
        tf = content.text_frame
        for bullet in slide_content["bullets"]:
            p = tf.add_paragraph()
            p.text = bullet

        # Add takeaway message at the bottom left with larger font
        left = Inches(1)
        top = Inches(6)
        width = Inches(6)
        height = Inches(2)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = slide_content["takeaway_message"]
        for paragraph in tf.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(20)  # Larger font size

        # Add notes to the slide
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        for talking_point in slide_content["talking_points"]:
            notes_tf.add_paragraph().text = talking_point

    # Save the presentation
    presentation.save("SlideDeck.pptx")

def get_download_link(file_path):
    with open(file_path, 'rb') as f:
        data = f.read()
    b64 = base64.b64encode(data).decode()
    return f'<a href="data:application/octet-stream;base64,{b64}" download="{file_path}">Download file</a>'

def reset_all():
    for key in list(st.session_state.keys()):
        del st.session_state[key]
    st.session_state.confirm_details = False  # Reset the checkbox value
     
def format_slide_content(slide_content):
    formatted_content = ""
    for key, value in slide_content.items():
        formatted_content += f"{key.capitalize()}:\n"  # Add ':' here if you want to separate the key from the value with a colon
        if isinstance(value, Iterable) and not isinstance(value, str):
            for i, item in enumerate(value, start=1):
                formatted_content += f"  {i}. {item}\n"
        else:
            formatted_content += f"  {value}\n"  # Add '\n' here to start value from a new line
        formatted_content += "\n"  # Add '\n' here to create a space between different key-value pairs
    return formatted_content

def setup_app_title():
    st.markdown("""
    <style>
    .big-font {
        font-size:50px !important;
        color: purple;
        font-weight: bold;
    }
    </style>
    """, unsafe_allow_html=True)
    st.markdown('<p class="big-font">🎨 SlideSage: Crafting <span style="color: teal;">Powerful</span> Presentations with <span style="color: pink;">AI</span> 🚀</p>', unsafe_allow_html=True)

def setup_sidebar_style():
    st.markdown("""
    <style>
    .reportview-container .main .block-container {
        margin-left: 10px;
        margin-right: 10px;
        padding-top: 10px;
        padding-right: 10px;
        padding-left: 10px;
        padding-bottom: 10px;
    }
    .sidebar .sidebar-content {
        background-color: #f0f0f5;
    }
    </style>
    """, unsafe_allow_html=True)

def main():
    setup_app_title()
    setup_sidebar_style()
    # Step 1: Allow user to enter a topic
    presentation_topic = st.sidebar.text_input('Topic for the PowerPoint Deck')

    # Step 2: Allow the user to select n charts for the outline
    num_slides = st.sidebar.number_input('Number of slides', min_value=1)
    
    # Step 7: Prompt the user to enter their presenter name, presentation title, and company name
    st.sidebar.title('Presentation Details')
    company_name = st.sidebar.text_input('Company name', 'Company')
    presentation_name = st.sidebar.text_input('Presentation name', 'Presentation')
    presenter = st.sidebar.text_input('Presenter', 'Presenter')

    # Step 2.1: Allow user to select engine
    engine = st.sidebar.selectbox('Select model', ['gpt-3.5-turbo', 'gpt-4'])

    # Step 3: Show estimated tokens and cost
    estimated_tokens = num_slides * TOKENS_PER_SLIDE_ESTIMATE
    st.sidebar.write(f"Estimated token usage: {estimated_tokens}")
    cost_per_1k_tokens = 0.002 if engine == 'gpt-3.5-turbo' else 0.02  # Adjust cost based on engine
    estimated_cost = estimated_tokens / 1000 * cost_per_1k_tokens
    st.sidebar.write(f"Estimated cost: ${estimated_cost}")

    # Before the 'Generate Outline' button press
    if 'outline' not in st.session_state:
        st.session_state['outline'] = []
        st.session_state['approved'] = False

    if estimated_tokens > MAX_TOKENS:
        st.warning(f"Estimated token usage is {estimated_tokens}, which is more than the maximum allowed ({MAX_TOKENS}). Consider reducing the number of slides.")
    else:
        
        # Step 4: Generate the outline upon pressing Generate Outline
        if st.sidebar.button('Generate Outline'):
            try:
                st.session_state['outline'] = generate_outline(presentation_topic, num_slides, engine)
                st.session_state['outline_copy'] = list(st.session_state['outline'])  # Create a copy of the outline

                # Step 5: Display the editable outline in the sidebar
                if st.session_state['outline']:
                    st.sidebar.write('Generated Outline:')
                    for i, slide_title in enumerate(st.session_state['outline_copy']):
                        st.session_state['outline_copy'][i] = st.sidebar.text_input(f'Slide {i+1}', slide_title)

                    # Add a button for the user to confirm their edits
                    if st.sidebar.button('Confirm Outline Edits'):
                        st.session_state['outline'] = list(st.session_state['outline_copy'])  # Update the outline
                        st.session_state['outline_edited'] = True
            except Exception as e:
                st.error(f"Failed to generate outline: {e}")


        # Step 9: Generate slide content for each slide title in the outline
        if 'outline_edited' in st.session_state and st.session_state['outline_edited']:
            slides_content = []
            for slide_title in st.session_state['outline']:
                st.write(f"Generating slide content for: {slide_title}")
                slide_content = generate_slide_content(slide_title, engine)
                slides_content.append(slide_content)

            st.session_state['slides_content'] = slides_content


    # Step 10: Show the "Create Presentation" button 
    if 'slides_content' in st.session_state and st.session_state['slides_content']:
        if st.sidebar.button('Create Presentation'):
            
            # Display the slide content in a formatted manner
            #st.write(f"Slide Content:\n{format_slide_content(slide_content)}")
            
            create_presentation(st.session_state['slides_content'], company_name, presentation_name, presenter)
            st.success('Presentation created successfully!')
            
            # Step 11: Allow the user to download the presentation with a link
            download_link = get_download_link("SlideDeck.pptx")
            st.markdown(download_link, unsafe_allow_html=True)

            
     # Reset Button
    if st.sidebar.button('Reset'):
        reset_all()
      
if __name__ == "__main__":
    main()



